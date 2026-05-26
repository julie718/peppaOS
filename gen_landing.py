"""Convert LumiOS Markdown proposals to PPTX + HTML (for PDF) — v3"""
import re, os

ROOT = os.path.dirname(os.path.abspath(__file__))
DOCS = os.path.join(ROOT, "docs")

FILES = [
    ("LumiOS_ToB_企业方案.md", "ToB"),
    ("LumiOS_ToC_消费方案.md", "ToC"),
    ("LumiOS_落地路径方案.md", "落地路径"),
]

# ── palette ──
C = {
    'ink':       (0x10, 0x10, 0x28),
    'white':     (0xFF, 0xFF, 0xFF),
    'accent':    (0x4F, 0x46, 0xE5),  # indigo
    'accent2':   (0x7C, 0x3A, 0xED),  # violet
    'body':      (0x33, 0x33, 0x44),
    'muted':     (0x78, 0x78, 0x90),
    'light':     (0xF3, 0xF4, 0xFA),
    'card':      (0xF8, 0xF8, 0xFF),
    'gold':      (0xF5, 0x9E, 0x0B),
    'green':     (0x10, 0xB9, 0x81),
    'red':       (0xEF, 0x44, 0x44),
    'table_hdr': (0xEE, 0xF0, 0xFF),
    'table_alt': (0xFA, 0xFA, 0xFF),
    'callout':   (0xEE, 0xF0, 0xFF),
}

def parse_md(filepath):
    with open(filepath, "r", encoding="utf-8") as f:
        text = f.read()
    slides_raw = re.split(r'\n## 第\d+页 · ', text)
    h1_match = re.search(r'^# (.+)$', text, re.MULTILINE)
    cover_title = h1_match.group(1).strip() if h1_match else "LumiOS"
    slides = []
    for s in slides_raw[1:]:
        parts = s.strip().split('\n', 1)
        title = parts[0].strip()
        body = parts[1].strip() if len(parts) > 1 else ""
        slides.append((title, body))
    return cover_title, slides

def parse_body(body_text):
    """Parse body into structured blocks: paragraphs, tables, h3s, callouts, lists, code"""
    blocks = []
    lines = body_text.split('\n')
    i = 0
    while i < len(lines):
        clean = lines[i].strip()
        if not clean:
            i += 1; continue

        # table
        if clean.startswith('|'):
            rows = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                cells = [c.strip() for c in lines[i].strip().split('|') if c.strip()]
                if not all(c.replace('-','').replace(':','') == '' for c in cells):
                    rows.append([c.strip('*') for c in cells])
                i += 1
            blocks.append(('table', rows))
            continue

        # code block (ascii art)
        if '─' in clean or '┌' in clean or '├' in clean or '└' in clean or '│' in clean or '↓' in clean:
            code_lines = []
            while i < len(lines) and lines[i].strip():
                code_lines.append(lines[i].strip())
                i += 1
            blocks.append(('code', '\n'.join(code_lines)))
            continue

        # h3
        if clean.startswith('###'):
            blocks.append(('h3', clean[3:].strip()))
            i += 1; continue

        # bold callout (standalone **...**) — only if entire line is bold
        if clean.startswith('**') and clean.endswith('**') and len(clean) < 150:
            blocks.append(('callout', clean.strip('*')))
            i += 1; continue

        # list items (consecutive)
        if clean.startswith('- '):
            items = []
            while i < len(lines) and lines[i].strip().startswith('- '):
                items.append(lines[i].strip()[2:])
                i += 1
            blocks.append(('list', items))
            continue

        # regular paragraph
        blocks.append(('p', clean))
        i += 1

    return blocks

# ═══════════════════════════════════════════════════════════
#  HTML / PDF  —  tight magazine layout
# ═══════════════════════════════════════════════════════════

def fmt_inline(text):
    """Convert **bold** and handle inline formatting"""
    return re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', text)

def gen_html(cover_title, slides, outpath):
    css = """<!DOCTYPE html><html lang="zh-CN"><head><meta charset="UTF-8">
<style>
  @page { size: A4; margin: 10mm 12mm; }
  * { margin:0; padding:0; box-sizing:border-box; }
  body { font-family:'Microsoft YaHei','PingFang SC',sans-serif; color:#1a1a2e; line-height:1.7; background:#fff; }
  .cover {
    text-align:center; page-break-after:always;
    background:linear-gradient(160deg,#05051a 0%,#0f0c29 30%,#1a1040 60%,#0d0b1e 100%);
    color:#fff; position:relative; overflow:hidden;
    display:flex; flex-direction:column; justify-content:center;
    min-height:100vh; padding:8% 10%;
  }
  .cover::before {
    content:""; position:absolute; top:-30%; right:-15%;
    width:70%; height:90%;
    background:radial-gradient(circle,rgba(99,102,241,.12) 0%,transparent 70%);
    border-radius:50%;
  }
  .cover::after {
    content:""; position:absolute; bottom:-20%; left:-10%;
    width:50%; height:60%;
    background:radial-gradient(circle,rgba(139,92,246,.08) 0%,transparent 70%);
    border-radius:50%;
  }
  .cover>* { position:relative; z-index:1; }
  .cover .tag {
    display:inline-block; border:1.5pt solid rgba(129,140,248,.4);
    color:#a5b4fc; font-size:9pt; letter-spacing:3pt; padding:4pt 16pt;
    border-radius:20pt; margin-bottom:24pt;
  }
  .cover h1 {
    font-size:40pt; letter-spacing:6pt; margin-bottom:10pt;
    font-weight:900; text-shadow:0 2px 20px rgba(99,102,241,.3);
  }
  .cover .deck {
    font-size:14pt; color:#c7d2fe; line-height:2;
    margin-bottom:28pt; letter-spacing:1pt;
  }
  .cover .firm {
    font-size:10pt; color:#6b7280; letter-spacing:2pt;
    border-top:1pt solid rgba(255,255,255,.1);
    padding-top:16pt; width:60%; margin:0 auto;
  }
  .cover .accent-line {
    width:60pt; height:2pt; background:#818cf8; margin:0 auto 24pt auto;
    border-radius:1pt;
  }
  .section { page-break-after:always; padding:0; }
  .section h2 { font-size:20pt; color:#3730a3; font-weight:800; border-left:5pt solid #818cf8; padding:4pt 0 4pt 12pt; margin-bottom:14pt; line-height:1.2; }
  .section h3 { font-size:12pt; color:#4338ca; font-weight:700; margin:14pt 0 6pt 0; }
  .section h3::before { content:"▸ "; color:#818cf8; }
  .section p { font-size:11pt; color:#334; margin-bottom:7pt; text-align:justify; }
  .section strong { color:#1e1b4b; }
  .section li { font-size:11pt; color:#334; margin-bottom:4pt; }
  .section ul { padding-left:16pt; margin:6pt 0; }
  .callout { background:#eef2ff; border-left:4pt solid #6366f1; padding:8pt 14pt; margin:12pt 0; border-radius:0 5pt 5pt 0; font-weight:600; color:#312e81; font-size:11pt; }
  .section table { border-collapse:collapse; width:100%; font-size:9.5pt; margin:10pt 0; }
  .section th { background:#e0e4ff; color:#3730a3; font-weight:700; padding:6pt 9pt; text-align:left; font-size:9.5pt; }
  .section td { border-bottom:1px solid #e8e8f0; padding:5pt 9pt; color:#445; font-size:9.5pt; }
  .section tr:last-child td { border-bottom:none; }
  .section pre { background:#f4f4fa; padding:8pt 12pt; border-radius:4pt; font-family:'Consolas','Courier New',monospace; font-size:8.5pt; line-height:1.45; color:#3730a3; overflow-x:auto; margin:8pt 0; white-space:pre-wrap; }
  hr { border:none; border-top:1px solid #d8d8e8; margin:10pt 0; }
</style></head><body>"""

    with open(outpath, "w", encoding="utf-8") as f:
        f.write(css)
        # Build cover subtitle from slides content
        cover_deck = ""
        for title, _ in slides[:1]:
            if title == '封面':
                # Extract the cover description
                pass
        f.write(f'<div class="cover">'
                f'<div class="tag">AI 人格操作系统</div>'
                f'<h1>{cover_title}</h1>'
                f'<div class="accent-line"></div>'
                f'<div class="deck">让每一个组织，长出属于自己的 AI。</div>'
                f'<div class="firm">浙江灵序科技有限公司 &nbsp;|&nbsp; lumiai.asia &nbsp;|&nbsp; 2026</div>'
                f'</div>\n')

        for title, body_text in slides:
            if title == '封面':
                continue
            blocks = parse_body(body_text)
            f.write(f'<div class="section"><h2>{title}</h2>\n')

            for kind, data in blocks:
                if kind == 'table':
                    f.write('<table>\n')
                    for ri, row in enumerate(data):
                        tag = 'th' if ri == 0 else 'td'
                        f.write('<tr>' + ''.join(f'<{tag}>{fmt_inline(c)}</{tag}>' for c in row) + '</tr>\n')
                    f.write('</table>\n')
                elif kind == 'h3':
                    f.write(f'<h3>{fmt_inline(data)}</h3>\n')
                elif kind == 'callout':
                    f.write(f'<div class="callout">{fmt_inline(data)}</div>\n')
                elif kind == 'bold':
                    f.write(f'<p><strong>{fmt_inline(data)}</strong></p>\n')
                elif kind == 'list':
                    f.write('<ul>' + ''.join(f'<li>{fmt_inline(it)}</li>' for it in data) + '</ul>\n')
                elif kind == 'code':
                    f.write(f'<pre>{data}</pre>\n')
                elif kind == 'p':
                    if data == '---':
                        continue  # skip markdown horizontal rules
                    else:
                        f.write(f'<p>{fmt_inline(data)}</p>\n')

            f.write('</div>\n')
        f.write('</body></html>')
    print(f"  HTML: {outpath}")

# ═══════════════════════════════════════════════════════════
#  PPTX  —  clean modern design with visual elements
# ═══════════════════════════════════════════════════════════

def gen_pptx(cover_title, slides, outpath, label):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn

    prs = Presentation()
    W, H = Inches(13.333), Inches(7.5)
    prs.slide_width, prs.slide_height = W, H

    def rgb(c): return RGBColor(*c)

    def bg(slide, color):
        f = slide.background.fill; f.solid(); f.fore_color.rgb = rgb(color)

    def rect(slide, l, t, w, h, color, radius=None):
        """Add a filled rectangle shape, optionally rounded"""
        import pptx.enum.shapes as sh
        shape = slide.shapes.add_shape(
            sh.MSO_SHAPE.ROUNDED_RECTANGLE if radius else sh.MSO_SHAPE.RECTANGLE,
            Inches(l), Inches(t), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
        return shape

    def circle(slide, l, t, d, color):
        shape = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(d), Inches(d))  # oval
        shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color)
        shape.line.fill.background()
        return shape

    def tb(slide, l, t, w, h, text, size=Pt(11), color=C['body'], bold=False,
           align=PP_ALIGN.LEFT, font='Microsoft YaHei', anchor='t'):
        txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = txBox.text_frame; tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = size
        tf.paragraphs[0].font.bold = bold
        tf.paragraphs[0].font.color.rgb = rgb(color)
        tf.paragraphs[0].alignment = align
        tf.paragraphs[0].font.name = font
        return tf

    def mp(tf, text, size=Pt(11), color=C['body'], bold=False, align=PP_ALIGN.LEFT, before=Pt(0), after=Pt(1)):
        p = tf.add_paragraph()
        p.text = text; p.font.size = size; p.font.bold = bold
        p.font.color.rgb = rgb(color); p.alignment = align
        p.font.name = 'Microsoft YaHei'
        p.space_before = before; p.space_after = after
        return p

    def accent_bar(slide, y=0, h=0.04):
        bar = slide.shapes.add_shape(1, Inches(0), Inches(y), W, Inches(h))
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(C['accent']); bar.line.fill.background()

    def make_table(slide, data, l, t, w, row_h=0.32):
        """Create a styled table. data[0] is header."""
        rows, cols = len(data), max(len(r) for r in data)
        h = row_h * rows
        tbl_shape = slide.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h))
        tbl = tbl_shape.table
        cw = w / cols
        for ci in range(cols):
            tbl.columns[ci].width = Inches(cw)
        for ri in range(rows):
            for ci in range(min(cols, len(data[ri]))):
                cell = tbl.cell(ri, ci)
                cell.text = ""
                p = cell.text_frame.paragraphs[0]
                p.text = str(data[ri][ci])
                p.font.name = 'Microsoft YaHei'
                p.font.size = Pt(10) if ri == 0 else Pt(9)
                p.font.bold = (ri == 0)
                p.font.color.rgb = rgb(C['ink']) if ri == 0 else rgb(C['body'])
                # cell fill
                tc = cell._tc.get_or_add_tcPr()
                sf = cell._tc.makeelement(qn('a:solidFill'), {})
                clr_hex = 'EEF0FF' if ri == 0 else ('FAFAFF' if ri % 2 == 0 else 'FFFFFF')
                sc = cell._tc.makeelement(qn('a:srgbClr'), {'val': clr_hex})
                sf.append(sc); tc.append(sf)
                tc.set(qn('a:marL'), str(Emu(Inches(0.08))))
                tc.set(qn('a:marR'), str(Emu(Inches(0.08))))
                tc.set(qn('a:marT'), str(Emu(Inches(0.04))))
                tc.set(qn('a:marB'), str(Emu(Inches(0.04))))
        return tbl

    def big_number(slide, l, t, number, label_str, accent_color=C['accent']):
        """Large number with label below"""
        tf = tb(slide, l, t, 2.0, 1.2, str(number), Pt(36), accent_color, True, PP_ALIGN.CENTER)
        mp(tf, label_str, Pt(10), C['muted'], False, PP_ALIGN.CENTER, Pt(0))

    # ─── cover ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, C['ink'])
    # decorative top bar
    rect(slide, 0, 0, 13.333, 0.06, C['accent'])
    # large accent circle (decorative)
    circle(slide, 10.2, -0.5, 4.5, (0x1E, 0x1B, 0x4B))
    circle(slide, 9.8, 4.0, 3.0, (0x25, 0x20, 0x55))
    # title
    tf = tb(slide, 1.2, 1.8, 10.5, 1.8, cover_title, Pt(46), C['white'], True, PP_ALIGN.CENTER)
    mp(tf, f"{label}方案", Pt(22), C['accent'], True, PP_ALIGN.CENTER, Pt(12))
    # separator
    sep = slide.shapes.add_shape(1, Inches(4.5), Inches(4.2), Inches(4.333), Inches(0.025))
    sep.fill.solid(); sep.fill.fore_color.rgb = rgb(C['accent']); sep.line.fill.background()
    mp(tf, "浙江灵序科技有限公司  |  lumiai.asia  |  2026", Pt(12), C['muted'], False, PP_ALIGN.CENTER, Pt(16))

    # ─── body slides ───
    for pg, (title, body_text) in enumerate(slides):
        blocks = parse_body(body_text)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg(slide, C['white'])
        accent_bar(slide, 0, 0.05)

        # title
        tb(slide, 0.7, 0.35, 11.8, 0.7, title, Pt(26), C['ink'], True)
        # separator
        sep = slide.shapes.add_shape(1, Inches(0.7), Inches(1.05), Inches(3.0), Inches(0.025))
        sep.fill.solid(); sep.fill.fore_color.rgb = rgb(C['accent']); sep.line.fill.background()

        # page number
        tb(slide, 12.3, 7.05, 0.8, 0.35, str(pg + 2), Pt(9), C['muted'], False, PP_ALIGN.RIGHT)

        # layout content blocks
        y = 1.35
        L = 0.7
        W_content = 11.9

        for kind, data in blocks:
            if y > 6.8:
                break

            if kind == 'h3':
                # section header with small accent bar
                rect(slide, L, y + 0.08, 0.06, 0.28, C['accent'])
                tb(slide, L + 0.18, y, W_content, 0.38, data, Pt(15), C['accent'], True)
                y += 0.42

            elif kind == 'callout':
                rect(slide, L, y, 11.9, 0.52, C['callout'], radius=0.06)
                rect(slide, L, y, 0.05, 0.52, C['accent'])
                tb(slide, L + 0.25, y + 0.08, 11.4, 0.38, data, Pt(13), C['ink'], True)
                y += 0.62

            elif kind == 'bold':
                tb(slide, L, y, W_content, 0.32, data, Pt(14), C['ink'], True)
                y += 0.38

            elif kind == 'table':
                rh = 0.3
                th = min(rh * len(data), 6.8 - y)
                make_table(slide, data, L, y, W_content, rh)
                y += th + 0.2

            elif kind == 'list':
                for item in data[:12]:  # max 12 items per slide
                    if y > 6.8: break
                    # bullet dot
                    circle(slide, L + 0.05, y + 0.08, 0.1, C['accent'])
                    tb(slide, L + 0.3, y, 11.3, 0.26, item, Pt(10.5), C['body'])
                    y += 0.26
                y += 0.06

            elif kind == 'code':
                n_lines = data.count('\n') + 1
                code_h = min(n_lines * 0.22, 2.0)
                rect(slide, L, y, W_content, code_h + 0.16, (0xF5, 0xF5, 0xFC), radius=0.04)
                tb(slide, L + 0.2, y + 0.08, W_content - 0.4, code_h, data, Pt(7.5), C['accent'])
                y += code_h + 0.25

            elif kind == 'p':
                # detect if it looks like a key-value line or has special patterns
                text = data
                # check for number highlights: "目标：..." or "xxx 家" etc
                tb(slide, L, y, W_content, 0.28, text, Pt(10.5), C['body'])
                y += 0.30

        # bottom accent line
        rect(slide, 0, 7.47, 13.333, 0.03, C['accent'])

    prs.save(outpath)
    print(f"  PPTX: {outpath}")

# ─── Main ───
for filename, label in FILES:
    md_path = os.path.join(DOCS, filename)
    cover, slides = parse_md(md_path)
    base = os.path.splitext(filename)[0]
    print(f"\n{'='*40}\n{label}: {len(slides)} slides + cover")
    gen_html(cover, slides, os.path.join(DOCS, f"{base}.html"))
    gen_pptx(cover, slides, os.path.join(DOCS, f"{base}.pptx"), label)

print("\nDone. Run Edge headless to convert HTML → PDF:")
for filename, _ in FILES:
    base = os.path.splitext(filename)[0]
    html = os.path.join(DOCS, f"{base}.html")
    pdf = os.path.join(DOCS, f"{base}.pdf")
    print(f'  msedge --headless --print-to-pdf="{pdf}" "file:///{html}"')
